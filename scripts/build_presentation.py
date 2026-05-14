"""Generate the project overview PPTX deck.

Run: python scripts/build_presentation.py
Output: docs/MGT-AI-Solution-Overview.pptx
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = pathlib.Path(__file__).resolve().parents[1] / "docs" / "MGT-AI-Solution-Overview.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Brand palette
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
TEAL = RGBColor(0x0E, 0x8A, 0x8A)
AMBER = RGBColor(0xE3, 0x9E, 0x1C)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_band(slide, color=NAVY, height=Inches(0.55), top=0):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, height
    )
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    return band


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=NAVY, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return box


def add_bullets(slide, left, top, width, height, items, *,
                size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return box


def slide_header(slide, title, subtitle=None):
    add_band(slide, NAVY, Inches(0.9))
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GREY)


def add_box(slide, left, top, width, height, title, body, *,
            accent=TEAL):
    # accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    # card body
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(0.12), top, width - Inches(0.12), height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = LIGHT
    # title
    add_text(slide, left + Inches(0.3), top + Inches(0.15),
             width - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    # body
    tb = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.6),
        width - Inches(0.4), height - Inches(0.7),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = GREY
        run.font.name = "Segoe UI"
        p.space_after = Pt(3)


def add_node(slide, left, top, width, height, label, *,
             color=TEAL, text_color=WHITE):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Segoe UI"
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=GREY):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(1.75)


# ---------------------------------------------------------------- Slide 1: Title
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# Accent block
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(0.4), Inches(2.7))
acc.fill.solid(); acc.fill.fore_color.rgb = AMBER; acc.line.fill.background()

add_text(s, Inches(0.9), Inches(2.3), Inches(11), Inches(0.6),
         "MGT AI Solution", size=18, bold=True, color=AMBER)
add_text(s, Inches(0.9), Inches(2.85), Inches(12), Inches(1.2),
         "Agent-Driven Contract Extraction",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(4.05), Inches(12), Inches(0.6),
         "LangGraph • LLM Decisioning • Confidence Fusion • Dead-Letter Replay",
         size=18, color=LIGHT)

# Author block
add_text(s, Inches(0.9), Inches(5.4), Inches(12), Inches(0.45),
         "Done by", size=12, color=LIGHT)
add_text(s, Inches(0.9), Inches(5.75), Inches(12), Inches(0.55),
         "Shibam Samaddar", size=22, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(6.25), Inches(12), Inches(0.4),
         "AI Solutions Engineer", size=14, color=AMBER)

add_text(s, Inches(0.9), Inches(6.95), Inches(12), Inches(0.4),
         "Architecture & Agentic Flow Overview",
         size=12, color=LIGHT)

# ---------------------------------------------------------------- Slide 2: Problem & Why
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Problem & Why Scenario 2",
             "Why contract extraction exercises the full agent stack")

add_box(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5),
        "The Problem",
        [
            "Contracts arrive as raw text from many sources.",
            "Need structured fields: parties, dates, term, governing law.",
            "Some are clean, some malformed, some need a human.",
            "Must persist, audit, and recover from third-party API failures.",
        ], accent=AMBER)

add_box(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5),
        "Why Scenario 2",
        [
            "High-stakes routing: low confidence escalates to humans.",
            "Combines deterministic + LLM confidence (fusion).",
            "Schema-aware repair loops show genuine reasoning.",
            "Realistic tool-failure recovery with replay.",
        ], accent=TEAL)

add_box(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6),
        "Four Agentic Capabilities Demonstrated",
        [
            "1. Decisioning  — LangGraph routes branches based on the LLM's own outputs (not hard-coded rules).",
            "2. Tool use     — Agent decides whether to POST to the SOR; can also skip the tool entirely.",
            "3. Memory       — TypedDict state + SQLite + JSONL trace, all keyed by a single trace_id.",
            "4. Uncertainty  — Confidence fusion, evidence verification, schema repair, dead-letter queue.",
        ], accent=NAVY)

# ---------------------------------------------------------------- Slide 3: High-level Architecture
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "High-Level Architecture",
             "Caller → FastAPI → LangGraph → SQLite + Mock SOR")

# layered boxes
add_node(s, Inches(0.6), Inches(2.0), Inches(2.4), Inches(0.9),
         "HTTP Caller\n(smoke_http.py / cli.py)", color=GREY)
add_node(s, Inches(3.5), Inches(2.0), Inches(2.6), Inches(0.9),
         "FastAPI :8000\nsrc/api/main.py", color=NAVY)
add_node(s, Inches(6.6), Inches(2.0), Inches(3.0), Inches(0.9),
         "LangGraph StateGraph\n9 nodes • 3 cond. edges", color=TEAL)
add_node(s, Inches(10.2), Inches(2.0), Inches(2.5), Inches(0.9),
         "Mock SOR :8001\nsrc/api/sor.py", color=AMBER)

# arrows
add_arrow(s, Inches(3.0), Inches(2.45), Inches(3.5), Inches(2.45))
add_arrow(s, Inches(6.1), Inches(2.45), Inches(6.6), Inches(2.45))
add_arrow(s, Inches(9.6), Inches(2.45), Inches(10.2), Inches(2.45))

# Persistence layer
add_node(s, Inches(3.5), Inches(4.0), Inches(6.1), Inches(1.0),
         "SQLite Persistence\nruns • stage_events • artifacts • dead_letter", color=NAVY)
add_arrow(s, Inches(7.9), Inches(2.9), Inches(6.55), Inches(4.0))
add_arrow(s, Inches(11.4), Inches(2.9), Inches(8.5), Inches(4.0))

# LLM backends
add_node(s, Inches(0.6), Inches(4.0), Inches(2.4), Inches(1.0),
         "LLM Clients\nMock • Ollama • OpenAI", color=TEAL)
add_arrow(s, Inches(1.8), Inches(4.0), Inches(1.8), Inches(2.9))

# Observability
add_node(s, Inches(10.2), Inches(4.0), Inches(2.5), Inches(1.0),
         "JSONL Trace\nsolution/traces/{trace_id}", color=GREY)

add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5),
         "Single trace_id (UUID) is the primary key across HTTP responses, the database, "
         "JSONL traces, and the SOR webhook — enabling full audit & replay.",
         size=14, color=GREY)

# ---------------------------------------------------------------- Slide 4: LangGraph Node Flow
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Agent Flow — LangGraph Nodes",
             "9 nodes, 3 conditional edges; LLM outputs decide each branch")

# Layout: linear with branches
y0 = Inches(1.8)
add_node(s, Inches(0.4), y0, Inches(1.5), Inches(0.7), "ingest", color=GREY)
add_node(s, Inches(2.1), y0, Inches(1.6), Inches(0.7), "extract\n(LLM)", color=TEAL)
add_node(s, Inches(3.9), y0, Inches(1.6), Inches(0.7), "validate\n(schema)", color=TEAL)
add_node(s, Inches(5.7), y0, Inches(1.7), Inches(0.7), "route_decision", color=NAVY)
add_node(s, Inches(7.6), y0, Inches(1.7), Inches(0.7), "call_sor\n(retry x3)", color=AMBER)
add_node(s, Inches(9.5), y0, Inches(1.5), Inches(0.7), "persist", color=GREY)
add_node(s, Inches(11.2), y0, Inches(1.5), Inches(0.7), "end", color=GREY)

# arrows linear
for x1, x2 in [(1.9,2.1),(3.7,3.9),(5.5,5.7),(7.4,7.6),(9.3,9.5),(11.0,11.2)]:
    add_arrow(s, Inches(x1), Inches(2.15), Inches(x2), Inches(2.15))

# repair loop above validate
add_node(s, Inches(3.9), Inches(0.95), Inches(1.6), Inches(0.55),
         "repair_schema (loop ≤3)", color=AMBER)
add_arrow(s, Inches(4.7), Inches(1.8), Inches(4.7), Inches(1.5))
add_arrow(s, Inches(4.7), Inches(0.95), Inches(4.7), Inches(0.7))

# branches below route
add_node(s, Inches(5.7), Inches(3.4), Inches(1.7), Inches(0.7),
         "needs_review", color=AMBER)
add_node(s, Inches(5.7), Inches(4.4), Inches(1.7), Inches(0.7),
         "dead_letter", color=GREY)
add_arrow(s, Inches(6.55), Inches(2.5), Inches(6.55), Inches(3.4))
add_arrow(s, Inches(6.55), Inches(4.1), Inches(6.55), Inches(4.4))

# from call_sor failure -> dead_letter
add_arrow(s, Inches(8.45), Inches(2.5), Inches(7.4), Inches(4.65))
add_text(s, Inches(7.5), Inches(3.5), Inches(2), Inches(0.3),
         "SOR fail → DLQ", size=10, color=GREY)

# from review/dlq -> persist
add_arrow(s, Inches(7.4), Inches(3.7), Inches(9.5), Inches(2.5))
add_arrow(s, Inches(7.4), Inches(4.65), Inches(9.5), Inches(2.5))

# legend
add_text(s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
         "Conditional edges driven by LLM outputs:", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.4), [
    "validate  →  evidence_verified == False  →  repair_schema  (loop ≤3, else dead_letter)",
    "route     →  self_overall_confidence ≥ 0.7 AND evidence_verified  →  call_sor  (else needs_review / dead_letter)",
    "call_sor  →  POST /contracts succeeds  →  persist;   5xx after retries  →  dead_letter",
], size=13, color=GREY)

# ---------------------------------------------------------------- Slide 5: Confidence Fusion
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Confidence Fusion & Routing",
             "How the agent quantifies its own uncertainty")

add_box(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(2.4),
        "Two Confidence Sources",
        [
            "Deterministic (src/confidence.py): schema match score (0–1).",
            "    required fields present + no extras.",
            "LLM-based (graph.py extract/route): self_llm_confidence (0–1).",
            "    LLM emits its own belief in its extraction.",
        ], accent=TEAL)

add_box(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.4),
        "Fusion Formula",
        [
            "self_overall_confidence  =",
            "      0.6 × deterministic  +  0.4 × llm",
            "                               −  Σ penalties",
            "Penalties: −0.15 missing required, −0.10 unverified citation,",
            "                −0.05 per schema error.",
        ], accent=AMBER)

add_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
        "Routing Threshold (src/routing.py)",
        [
            "≥ 0.70  AND  evidence_verified == True   →   call_sor   (POST /contracts)",
            "<  0.70  OR  evidence_verified == False  →   needs_review   (escalate to human)",
            "Anomaly / repair loop exhausted / SOR fails after retries  →  dead_letter   (replayable)",
            "",
            "Evidence verification: LLM cites text substrings; agent confirms they actually appear in the source.",
        ], accent=NAVY)

# ---------------------------------------------------------------- Slide 6: Reliability & DLQ
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Reliability — Retries, Dead-Letter & Replay",
             "Tool calls are resilient; failures become first-class state")

add_box(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(2.6),
        "Retry Policy (tenacity)",
        [
            "3 total attempts.",
            "Exponential backoff: 1s → 2s → 4s.",
            "Trigger: HTTP 5xx only (4xx logged, no retry).",
            "Fallback: route to dead_letter after 3 fails.",
        ], accent=AMBER)

add_box(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.6),
        "Dead-Letter Queue",
        [
            "Table: dead_letter (trace_id PK).",
            "Stores reason, text_preview, replay_count.",
            "Replay returns a NEW trace_id (audit-preserving).",
            "Original DLQ row not deleted — full history kept.",
        ], accent=TEAL)

add_box(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
        "Replay Endpoint",
        [
            "POST /dead-letter/{trace_id}/replay   →   { new_trace_id, status, sor_id }",
            "Optional payload: { updated_model: 'llama3.1:8b-instruct' } to re-try with a different LLM.",
            "Use case: human fixes prompt or model; re-runs the entire trace; original anomaly preserved for audit.",
        ], accent=NAVY)

# ---------------------------------------------------------------- Slide 7: Memory & Observability
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Memory & Observability",
             "Three persistence surfaces, one trace_id")

add_box(s, Inches(0.5), Inches(1.7), Inches(4.1), Inches(5.3),
        "1. In-Graph State",
        [
            "AgentState TypedDict.",
            "Fields: text, extracted_fields,",
            "    confidence, evidence_verified,",
            "    review_reasons, sor_decision,",
            "    sor_id, status, artifacts.",
            "Lives only during a run.",
        ], accent=TEAL)

add_box(s, Inches(4.7), Inches(1.7), Inches(4.0), Inches(5.3),
        "2. SQLite (src/db.py)",
        [
            "runs           — trace metadata.",
            "stage_events   — per-node transitions.",
            "artifacts      — extracted JSON +",
            "                 sor_id + confidence.",
            "dead_letter    — failed docs +",
            "                 replay history.",
        ], accent=NAVY)

add_box(s, Inches(8.8), Inches(1.7), Inches(4.0), Inches(5.3),
        "3. JSONL Trace",
        [
            "solution/traces/{trace_id}.jsonl",
            "One line per node event:",
            "    {timestamp, stage, node,",
            "     fields, result, error}",
            "Human-readable audit log;",
            "diffable in code review.",
        ], accent=AMBER)

# ---------------------------------------------------------------- Slide 8: Test Coverage
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "Test Coverage — All 3 Branches Exercised",
             "tests/ + runs/ snapshots prove every path")

# table-ish
headers = ["Scenario", "Branch", "Confidence", "SOR ID", "Notes"]
rows = [
    ("s2_happy_msa",          "extract_result", "0.75+", "✓ posted", "High-quality contract; SOR accepts."),
    ("s2_happy_subscription", "extract_result", "0.72",  "✓ posted", "Partial fields; repair loop recovers."),
    ("s2_needsreview_nda",    "needs_review",   "0.45",  "—",        "NDA schema violation; escalates."),
    ("s2_deadletter_poison",  "dead_letter",    "0.02",  "—",        "Malformed JSON; SOR fails; queued for replay."),
]

table = s.shapes.add_table(len(rows) + 1, len(headers),
                           Inches(0.5), Inches(1.8),
                           Inches(12.3), Inches(3.3)).table

widths = [Inches(2.6), Inches(2.0), Inches(1.6), Inches(1.5), Inches(4.6)]
for i, w in enumerate(widths):
    table.columns[i].width = w

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(13)
            r.font.name = "Segoe UI"

for i, row in enumerate(rows, start=1):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.color.rgb = NAVY
                r.font.name = "Segoe UI"

add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
         "Run:  python -m unittest discover -s tests -v   →   8 tests, full graph branching coverage.",
         size=14, bold=True, color=NAVY)

add_bullets(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.5), [
    "tests/test_graph.py      — node execution & conditional edges (mocked LLM + SOR).",
    "tests/test_pipeline.py   — end-to-end across all 4 scenarios.",
    "runs/s2_*/                 — committed artefacts proving each branch on real text.",
], size=13, color=GREY)

# ---------------------------------------------------------------- Slide 9: LLM Backends & Demo
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_header(s, "LLM Backends & How to Demo",
             "Pluggable via the LLM Protocol  (src/llm/base.py)")

add_box(s, Inches(0.5), Inches(1.7), Inches(4.0), Inches(2.6),
        "Mock  (default in tests)",
        [
            "src/llm/mock_client.py",
            "Deterministic; zero setup.",
            "Used by every unit test.",
        ], accent=GREY)

add_box(s, Inches(4.7), Inches(1.7), Inches(4.0), Inches(2.6),
        "Ollama  (local, free)",
        [
            "src/llm/ollama_client.py",
            "ollama pull llama3.1:8b-instruct",
            "$env:AGENT_LLM='ollama'",
            "Offline, reproducible.",
        ], accent=TEAL)

add_box(s, Inches(8.9), Inches(1.7), Inches(4.0), Inches(2.6),
        "OpenAI  (hosted)",
        [
            "src/llm/openai_client.py",
            "$env:OPENAI_API_KEY='sk-…'",
            "$env:AGENT_LLM='openai'",
            "Production-grade quality.",
        ], accent=AMBER)

add_box(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
        "Demo Commands",
        [
            "python cli.py demo            →  run sample contract; print full trace.",
            "python cli.py runs            →  list all runs with confidence + sor_id.",
            "python cli.py dlq             →  show dead-letter queue.",
            "python cli.py replay <id>     →  re-run a DLQ trace; returns new trace_id.",
            "python smoke_http.py          →  end-to-end HTTP smoke test (3 terminals).",
        ], accent=NAVY)

# ---------------------------------------------------------------- Slide 10: Summary
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
         "Summary", size=32, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "What makes this 'agentic', not just an LLM call",
         size=16, color=AMBER)

points = [
    ("Decisioning",
     "LangGraph routes branches based on the LLM's own outputs — not hard-coded rules."),
    ("Tool Use",
     "Agent decides when to POST to SOR, retries on 5xx, and can choose to skip the tool."),
    ("Memory",
     "TypedDict state + SQLite + JSONL trace, all keyed by a single trace_id."),
    ("Uncertainty",
     "Confidence fusion (0.6 det + 0.4 LLM − penalties); evidence verification; schema repair."),
    ("Recovery",
     "Dead-letter queue + replay endpoint preserve audit history while enabling fixes."),
    ("Observability",
     "Three persistence surfaces; every run reproducible from trace_id."),
]

top = Inches(2.0)
row_h = Inches(0.78)
for i, (k, v) in enumerate(points):
    y = top + row_h * i
    # bullet dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.2),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = AMBER
    dot.line.fill.background()
    add_text(s, Inches(1.05), y, Inches(2.6), Inches(0.6),
             k, size=18, bold=True, color=AMBER)
    add_text(s, Inches(3.7), y + Inches(0.04), Inches(9.4), Inches(0.6),
             v, size=15, color=WHITE)

add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         "MGT AI Solution  •  Scenario 2  •  Contract Extraction Agent   "
         "—   Done by Shibam Samaddar, AI Solutions Engineer",
         size=12, color=LIGHT)

prs.save(OUT)
print(f"Wrote {OUT}")
